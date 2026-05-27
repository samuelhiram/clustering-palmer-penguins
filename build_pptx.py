"""Genera presentacion.pptx (5 slides) para el proyecto de Mineria de Datos."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2E, 0x4F)
TEAL = RGBColor(0x1F, 0x77, 0xB4)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

blank = prs.slide_layouts[6]


def add_bg_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    return bar


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_bullets(slide, left, top, width, height, items, size=18, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def slide_header(slide, title, subtitle=None):
    add_bg_bar(slide)
    add_text(slide, Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.5),
             title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.5),
                 subtitle, size=14, color=GRAY)


# ===== Slide 1: Portada =====
s = prs.slides.add_slide(blank)
# fondo
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.2),
         "Clustering de Palmer Penguins", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.8),
         "Comparación de K-Means, DBSCAN y Agglomerative Clustering",
         size=22, color=RGBColor(0xCC, 0xDD, 0xEE))
add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
         "Proyecto de Minería de Datos — 26 de mayo de 2026",
         size=18, color=RGBColor(0xAA, 0xBB, 0xCC))
add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
         "Práctica de Ciencia 2.0 · datos abiertos · código reproducible",
         size=14, color=RGBColor(0x99, 0xAA, 0xBB))

# ===== Slide 2: Problema y objetivo =====
s = prs.slides.add_slide(blank)
slide_header(s, "1. Problema y objetivo")
add_text(s, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.5),
         "Problema", size=22, bold=True, color=TEAL)
add_text(s, Inches(0.4), Inches(1.55), Inches(12.5), Inches(1.3),
         "¿Se pueden descubrir automáticamente los grupos naturales de "
         "pingüinos (especies) a partir únicamente de sus características "
         "físicas, sin usar la etiqueta de especie durante el entrenamiento?",
         size=18, color=DARK)

add_text(s, Inches(0.4), Inches(3.1), Inches(12.5), Inches(0.5),
         "Objetivo", size=22, bold=True, color=TEAL)
add_bullets(s, Inches(0.6), Inches(3.65), Inches(12), Inches(2.5), [
    "Comparar 3 algoritmos de clustering no supervisado: K-Means, DBSCAN y Agglomerative.",
    "Evaluar con métricas interna (Silhouette) y externa (Adjusted Rand Index vs. especie).",
    "Demostrar un flujo completo de minería de datos siguiendo Ciencia 2.0.",
], size=18)

add_text(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5),
         "Reproducibilidad: RANDOM_STATE = 42 · semillas fijadas · requirements.txt incluido",
         size=14, color=GRAY)

# ===== Slide 3: Dataset y variables =====
s = prs.slides.add_slide(blank)
slide_header(s, "2. Dataset y variables", "Palmer Penguins (Gorman et al., 2014)")

add_bullets(s, Inches(0.4), Inches(1.3), Inches(6.3), Inches(4.5), [
    "344 observaciones · 3 especies",
    "Adelie (152) · Gentoo (124) · Chinstrap (68)",
    "3 islas: Biscoe, Dream, Torgersen",
    "Licencia CC-0 — dominio público",
    "Fuente: github.com/allisonhorst/palmerpenguins",
], size=18)

add_text(s, Inches(0.4), Inches(4.6), Inches(6.3), Inches(0.5),
         "Variables numéricas usadas", size=18, bold=True, color=TEAL)
add_bullets(s, Inches(0.4), Inches(5.15), Inches(6.3), Inches(2.0), [
    "bill_length_mm",
    "bill_depth_mm",
    "flipper_length_mm",
    "body_mass_g",
], size=16)

# imagen pairplot
s.shapes.add_picture("c:/Users/user/mineria/images/pairplot.png",
                     Inches(6.9), Inches(1.1), height=Inches(6.0))

# ===== Slide 4: Métodos =====
s = prs.slides.add_slide(blank)
slide_header(s, "3. Métodos utilizados",
             "Preprocesamiento: dropna() + StandardScaler (4 variables numéricas)")

# Tabla simple con tres columnas
from pptx.util import Inches as I
cols = [
    ("K-Means", "n_clusters = 3\nn_init = 10",
     "Minimiza la suma de distancias al centroide. Rápido y simple, asume clusters convexos."),
    ("DBSCAN", "eps = 0.6\nmin_samples = 5",
     "Basado en densidad. Detecta ruido (etiqueta −1). No requiere fijar k."),
    ("Agglomerative", "n_clusters = 3\nlinkage = ward",
     "Jerárquico aglomerativo. Une iterativamente clusters minimizando la varianza."),
]
col_w = Inches(4.0)
left0 = Inches(0.55)
top0 = Inches(1.4)
for i, (name, params, desc) in enumerate(cols):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             left0 + I(0.15) + col_w * i, top0,
                             col_w - I(0.3), Inches(5.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFA)
    box.line.color.rgb = TEAL
    add_text(s, left0 + I(0.3) + col_w * i, top0 + Inches(0.2),
             col_w - I(0.6), Inches(0.6), name, size=22, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, left0 + I(0.3) + col_w * i, top0 + Inches(1.0),
             col_w - I(0.6), Inches(1.4), params, size=14, color=TEAL,
             align=PP_ALIGN.CENTER)
    add_text(s, left0 + I(0.3) + col_w * i, top0 + Inches(2.3),
             col_w - I(0.6), Inches(3.0), desc, size=15, color=DARK)

# ===== Slide 5: Resultados =====
s = prs.slides.add_slide(blank)
slide_header(s, "4. Resultados",
             "Métricas y proyección PCA — Agglomerative obtiene el mejor ARI")

# Mini tabla de métricas
rows = [
    ("Algoritmo", "Silhouette", "ARI", "Ruido"),
    ("K-Means",        "0.447", "0.793", "0"),
    ("DBSCAN",         "0.557", "0.584", "33"),
    ("Agglomerative",  "0.454", "0.916", "0"),
]
table_left = Inches(0.4)
table_top = Inches(1.4)
table = s.shapes.add_table(len(rows), 4, table_left, table_top,
                           Inches(5.4), Inches(2.3)).table
for c in range(4):
    table.columns[c].width = Inches(1.35)
for r_idx, row in enumerate(rows):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.size = Pt(14)
        run.font.name = "Calibri"
        if r_idx == 0:
            run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            run.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0xEA, 0xF4, 0xEA)
                                        if row[0] == "Agglomerative"
                                        else WHITE)

add_text(s, Inches(0.4), Inches(4.0), Inches(5.6), Inches(0.5),
         "Observaciones", size=18, bold=True, color=TEAL)
add_bullets(s, Inches(0.4), Inches(4.5), Inches(5.6), Inches(3.0), [
    "Agglomerative recupera mejor las 3 especies (ARI = 0.92).",
    "Tras mapear cluster → especie por mayoría: exactitud 96.8 % (331/342).",
    "Mapeo: Cluster 0 → Adelie, Cluster 1 → Gentoo, Cluster 2 → Chinstrap.",
    "DBSCAN tiene mejor compacidad interna pero fusiona Adelie/Chinstrap.",
    "K-Means es muy competitivo y mucho más rápido.",
], size=15)

# imagen PCA
s.shapes.add_picture("c:/Users/user/mineria/images/clusters_pca.png",
                     Inches(6.2), Inches(1.4), width=Inches(7.0))

# ===== Slide 6: Conclusiones =====
s = prs.slides.add_slide(blank)
slide_header(s, "5. Conclusiones")

add_text(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.5),
         "¿Qué funcionó mejor?", size=20, bold=True, color=TEAL)
add_bullets(s, Inches(0.6), Inches(1.65), Inches(12), Inches(1.3), [
    "Agglomerative (ward) y K-Means recuperan la estructura biológica de las 3 especies.",
    "DBSCAN no es óptimo: el dataset no tiene densidades muy distintas.",
], size=16)

add_text(s, Inches(0.4), Inches(3.2), Inches(12.5), Inches(0.5),
         "Limitaciones", size=20, bold=True, color=TEAL)
add_bullets(s, Inches(0.6), Inches(3.75), Inches(12), Inches(1.3), [
    "Solo 4 variables numéricas y dataset pequeño (342 filas tras limpieza).",
    "DBSCAN es muy sensible al parámetro eps.",
], size=16)

add_text(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.5),
         "Trabajo futuro", size=20, bold=True, color=TEAL)
add_bullets(s, Inches(0.6), Inches(5.55), Inches(12), Inches(1.8), [
    "Incluir variables categóricas (sex, island) con encoding.",
    "Probar Gaussian Mixture Models y reducción de dimensionalidad con UMAP/t-SNE.",
    "Validar estabilidad de clusters con bootstrap.",
], size=16)

prs.save("c:/Users/user/mineria/presentacion.pptx")
print("presentacion.pptx generada con", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
